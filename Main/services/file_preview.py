"""Preview helpers for common local document formats.

The UI should not need to know how each document family is parsed.  This module
keeps preview extraction small, defensive, and dependency-isolated so more
preview handlers can be added later without bloating the main window code.
"""

from __future__ import annotations

import gzip
from html import escape
from pathlib import Path
import tarfile
from typing import Iterable
import zipfile

TEXT_PREVIEW_SUFFIXES = {
    ".txt", ".md", ".py", ".js", ".ts", ".tsx", ".html", ".htm", ".css",
    ".json", ".csv", ".xml", ".yml", ".yaml", ".toml", ".ini", ".log",
    ".sql", ".sh", ".ps1", ".bat", ".java", ".c", ".cpp", ".h", ".hpp",
    ".cs", ".go", ".rs", ".php", ".rb",
}
IMAGE_PREVIEW_SUFFIXES = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}
OFFICE_PREVIEW_SUFFIXES = {".docx", ".pptx", ".xlsx"}
PDF_PREVIEW_SUFFIXES = {".pdf"}
VIDEO_PREVIEW_SUFFIXES = {".mp4", ".mov", ".m4v", ".avi", ".mkv", ".webm"}
AUDIO_PREVIEW_SUFFIXES = {".mp3", ".wav", ".m4a", ".flac", ".ogg"}
ARCHIVE_PREVIEW_SUFFIXES = {".zip", ".7z", ".rar", ".tar", ".gz", ".tgz", ".tar.gz"}
PREVIEW_SUFFIXES = (
    TEXT_PREVIEW_SUFFIXES
    | IMAGE_PREVIEW_SUFFIXES
    | OFFICE_PREVIEW_SUFFIXES
    | PDF_PREVIEW_SUFFIXES
    | VIDEO_PREVIEW_SUFFIXES
    | AUDIO_PREVIEW_SUFFIXES
    | ARCHIVE_PREVIEW_SUFFIXES
)


_PREVIEW_CSS = """
<style>
body { margin: 0; }
.structured-preview {
    color: #dbeafe;
    font-family: "Inter", "Segoe UI", "SF Pro Display", sans-serif;
    font-size: 14px;
    line-height: 1.55;
}
.preview-hero {
    background: #111827;
    border: 1px solid #2f3b52;
    border-radius: 16px;
    padding: 16px 18px;
    margin-bottom: 14px;
}
.preview-kicker {
    color: #93c5fd;
    font-size: 12px;
    font-weight: 800;
    letter-spacing: 0.08em;
    text-transform: uppercase;
}
.preview-title {
    color: #f8fbff;
    font-size: 22px;
    font-weight: 850;
    margin-top: 4px;
}
.preview-meta {
    color: #9fb0c8;
    font-size: 13px;
    margin-top: 6px;
}
.preview-section {
    background: #0f172a;
    border: 1px solid #273449;
    border-radius: 14px;
    padding: 14px 16px;
    margin: 12px 0;
}
.preview-section-title {
    color: #bfdbfe;
    font-size: 15px;
    font-weight: 800;
    margin-bottom: 8px;
}
.preview-paragraph {
    color: #dbe7f7;
    margin: 7px 0;
}
.preview-slide-title,
.preview-sheet-title,
.preview-table-title {
    color: #ffffff;
    font-weight: 800;
    margin: 10px 0 6px 0;
}
ul.preview-list {
    margin-top: 6px;
    margin-bottom: 4px;
    padding-left: 20px;
}
.preview-list li {
    margin: 5px 0;
}
table.preview-table {
    width: 100%;
    border-collapse: collapse;
    margin-top: 8px;
    margin-bottom: 10px;
}
table.preview-table td,
table.preview-table th {
    border: 1px solid #28364b;
    padding: 6px 8px;
    color: #dbe7f7;
    vertical-align: top;
}
table.preview-table tr:nth-child(odd) td {
    background: #111c2f;
}
table.preview-table tr:nth-child(even) td {
    background: #0f172a;
}
.preview-muted {
    color: #93a4bc;
    font-size: 13px;
    font-style: italic;
}
.preview-document-pages {
    display: block;
}
.preview-page {
    background: #f8fafc;
    color: #172033;
    border: 1px solid #d8e2ef;
    border-radius: 10px;
    padding: 34px 38px;
    margin: 16px auto;
    max-width: 860px;
    min-height: 520px;
    box-shadow: 0 18px 40px rgba(0, 0, 0, 0.22);
}
.preview-page .preview-paragraph {
    color: #172033;
    font-family: "Cambria", "Georgia", serif;
    font-size: 16px;
    line-height: 1.65;
}
.preview-page .preview-heading {
    color: #0f2440;
    font-family: "Segoe UI", sans-serif;
    font-size: 22px;
    font-weight: 850;
    margin: 18px 0 8px 0;
}
.preview-page-number {
    color: #66758a;
    font-size: 12px;
    font-weight: 700;
    text-align: right;
    margin-top: 24px;
}
.preview-slide {
    background: #f8fafc;
    border: 1px solid #d8e2ef;
    border-radius: 12px;
    color: #172033;
    margin: 16px auto;
    max-width: 900px;
    min-height: 420px;
    padding: 28px 34px;
    box-shadow: 0 16px 36px rgba(0, 0, 0, 0.2);
}
.preview-slide .preview-slide-title {
    color: #0f2440;
    font-size: 20px;
}
.preview-slide ul.preview-list li {
    color: #172033;
    font-size: 16px;
    line-height: 1.45;
}
.preview-sheet {
    background: #f8fafc;
    border: 1px solid #d8e2ef;
    border-radius: 12px;
    color: #172033;
    margin: 16px 0;
    padding: 18px;
    overflow-x: auto;
}
.preview-sheet table.preview-table td,
.preview-sheet table.preview-table th {
    border: 1px solid #cbd8e8;
    color: #172033;
}
.preview-sheet table.preview-table tr:nth-child(odd) td {
    background: #ffffff;
}
.preview-sheet table.preview-table tr:nth-child(even) td {
    background: #eef4fb;
}
</style>
"""


def can_preview_with_handler(path: Path) -> bool:
    """Return True when this module has a structured preview handler."""
    return preview_kind(path) in {"office", "archive"}


def can_preview(path: Path) -> bool:
    """Return True when the app can show an inline preview surface."""
    return preview_kind(path) != "unsupported"


def preview_kind(path: Path) -> str:
    """Classify a local file into the preview surface it should use."""
    path = Path(path)
    suffix = _preview_suffix(path)

    if suffix in TEXT_PREVIEW_SUFFIXES:
        return "text"
    if suffix in IMAGE_PREVIEW_SUFFIXES:
        return "image"
    if suffix in OFFICE_PREVIEW_SUFFIXES:
        return "office"
    if suffix in PDF_PREVIEW_SUFFIXES:
        return "pdf"
    if suffix in VIDEO_PREVIEW_SUFFIXES:
        return "video"
    if suffix in AUDIO_PREVIEW_SUFFIXES:
        return "audio"
    if suffix in ARCHIVE_PREVIEW_SUFFIXES:
        return "archive"
    return "unsupported"


def _preview_suffix(path: Path) -> str:
    name = Path(path).name.lower()
    if name.endswith(".tar.gz"):
        return ".tar.gz"
    return Path(path).suffix.lower()


def structured_preview(path: Path, max_chars: int | None = None) -> str:
    """Return a readable plain-text preview for supported Office-style files."""
    path = Path(path)
    suffix = path.suffix.lower()

    try:
        if suffix == ".docx":
            return _truncate(_preview_docx_text(path), max_chars)
        if suffix == ".pptx":
            return _truncate(_preview_pptx_text(path), max_chars)
        if suffix == ".xlsx":
            return _truncate(_preview_xlsx_text(path), max_chars)
        if preview_kind(path) == "archive":
            return _truncate(_preview_archive_text(path), max_chars)
    except Exception as error:  # Defensive: previews should never crash UI.
        return f"Could not generate document preview:\n{error}"

    return ""


def structured_preview_html(path: Path, max_chars: int | None = None) -> str:
    """Return a styled HTML preview for supported Office-style files."""
    path = Path(path)
    suffix = path.suffix.lower()

    try:
        if suffix == ".docx":
            return _truncate_html(_preview_docx_html(path), max_chars)
        if suffix == ".pptx":
            return _truncate_html(_preview_pptx_html(path), max_chars)
        if suffix == ".xlsx":
            return _truncate_html(_preview_xlsx_html(path), max_chars)
        if preview_kind(path) == "archive":
            return _truncate_html(_preview_archive_html(path), max_chars)
    except Exception as error:
        return _html_page(
            "Preview unavailable",
            "Document preview",
            f"Could not generate document preview: {escape(str(error))}",
        )

    return ""


def is_structured_preview_html(text: str) -> bool:
    return "class=\"structured-preview\"" in (text or "")


def _html_page(kicker: str, title: str, body: str, meta: str = "") -> str:
    meta_html = f'<div class="preview-meta">{meta}</div>' if meta else ""
    return (
        _PREVIEW_CSS
        + '<div class="structured-preview">'
        + '<div class="preview-hero">'
        + f'<div class="preview-kicker">{escape(kicker)}</div>'
        + f'<div class="preview-title">{escape(title)}</div>'
        + meta_html
        + '</div>'
        + body
        + '</div>'
    )


def _preview_docx_html(path: Path) -> str:
    try:
        from docx import Document
    except Exception as error:
        return _html_page("Microsoft Word", path.name, f'<div class="preview-muted">DOCX preview requires python-docx. {escape(str(error))}</div>')

    document = Document(str(path))
    paragraphs = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    body_parts = []

    body_parts.append('<div class="preview-section"><div class="preview-section-title">Paragraphs</div>')
    if paragraphs:
        for paragraph in paragraphs[:120]:
            body_parts.append(f'<div class="preview-paragraph">{escape(paragraph)}</div>')
        if len(paragraphs) > 120:
            body_parts.append(f'<div class="preview-muted">Showing first 120 paragraphs out of {len(paragraphs)}.</div>')
    else:
        body_parts.append('<div class="preview-muted">No paragraph text found.</div>')
    body_parts.append('</div>')

    if document.tables:
        body_parts.append('<div class="preview-section"><div class="preview-section-title">Tables</div>')
        for table_index, table in enumerate(document.tables[:6], start=1):
            body_parts.append(f'<div class="preview-table-title">Table {table_index}</div>')
            body_parts.append('<table class="preview-table">')
            for row in table.rows[:16]:
                body_parts.append('<tr>')
                for cell in row.cells:
                    body_parts.append(f'<td>{escape(_normalise_cell(cell.text))}</td>')
                body_parts.append('</tr>')
            body_parts.append('</table>')
            if len(table.rows) > 16:
                body_parts.append(f'<div class="preview-muted">Showing first 16 rows out of {len(table.rows)}.</div>')
        if len(document.tables) > 6:
            body_parts.append(f'<div class="preview-muted">Showing first 6 tables out of {len(document.tables)}.</div>')
        body_parts.append('</div>')

    meta = f"{len(paragraphs)} paragraph(s) · {len(document.tables)} table(s)"
    return _html_page("Microsoft Word", path.name, "".join(body_parts), meta)


def _preview_pptx_html(path: Path) -> str:
    import re
    import zipfile
    import xml.etree.ElementTree as ET

    def slide_sort_key(name: str) -> int:
        match = re.search(r"slide(\d+)\.xml$", name)
        return int(match.group(1)) if match else 0

    with zipfile.ZipFile(path) as archive:
        slide_names = sorted(
            [name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")],
            key=slide_sort_key,
        )
        body_parts = []
        body_parts.append('<div class="preview-section"><div class="preview-section-title">Slides</div>')
        for slide_number, slide_name in enumerate(slide_names[:30], start=1):
            root = ET.fromstring(archive.read(slide_name))
            text_nodes = []
            for node in root.iter():
                if node.tag.endswith("}t") and node.text:
                    value = node.text.strip()
                    if value:
                        text_nodes.append(value)

            body_parts.append(f'<div class="preview-slide-title">Slide {slide_number}</div>')
            if text_nodes:
                body_parts.append('<ul class="preview-list">')
                for value in text_nodes:
                    body_parts.append(f'<li>{escape(value)}</li>')
                body_parts.append('</ul>')
            else:
                body_parts.append('<div class="preview-muted">No text found on this slide.</div>')
        if len(slide_names) > 30:
            body_parts.append(f'<div class="preview-muted">Showing first 30 slides out of {len(slide_names)}.</div>')
        body_parts.append('</div>')

    return _html_page("Microsoft PowerPoint", path.name, "".join(body_parts), f"{len(slide_names)} slide(s)")


def _preview_xlsx_html(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except Exception as error:
        return _html_page("Microsoft Excel", path.name, f'<div class="preview-muted">Excel preview requires openpyxl. {escape(str(error))}</div>')

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    try:
        body_parts = []
        body_parts.append('<div class="preview-section"><div class="preview-section-title">Workbook sheets</div>')
        for sheet_name in workbook.sheetnames[:8]:
            sheet = workbook[sheet_name]
            body_parts.append(f'<div class="preview-sheet-title">{escape(sheet.title)}</div>')
            rows = list(sheet.iter_rows(min_row=1, max_row=min(sheet.max_row or 1, 25), values_only=True))
            rendered_rows = [row for row in rows if any(cell is not None for cell in row)]
            if not rendered_rows:
                body_parts.append('<div class="preview-muted">No cell values found.</div>')
                continue
            body_parts.append('<table class="preview-table">')
            for row in rendered_rows:
                values = [_format_cell(value) for value in row]
                while values and values[-1] == "":
                    values.pop()
                body_parts.append('<tr>')
                for value in values:
                    body_parts.append(f'<td>{escape(value)}</td>')
                body_parts.append('</tr>')
            body_parts.append('</table>')
            if sheet.max_row and sheet.max_row > 25:
                body_parts.append(f'<div class="preview-muted">Showing first 25 rows out of {sheet.max_row}.</div>')
        if len(workbook.sheetnames) > 8:
            body_parts.append(f'<div class="preview-muted">Showing first 8 sheets out of {len(workbook.sheetnames)}.</div>')
        body_parts.append('</div>')
        return _html_page("Microsoft Excel", path.name, "".join(body_parts), f"{len(workbook.sheetnames)} sheet(s)")
    finally:
        workbook.close()


def _preview_docx_html(path: Path) -> str:
    """Render DOCX content as page-like HTML without artificial item limits."""
    try:
        from docx import Document
    except Exception as error:
        return _html_page("Microsoft Word", path.name, f'<div class="preview-muted">DOCX preview requires python-docx. {escape(str(error))}</div>')

    document = Document(str(path))
    blocks = list(_iter_docx_blocks(document))
    paragraph_count = sum(1 for block in blocks if block[0] == "paragraph" and block[1])
    table_count = sum(1 for block in blocks if block[0] == "table")

    body_parts = ['<div class="preview-document-pages"><div class="preview-page">']
    rendered_on_page = 0
    page_number = 1

    if blocks:
        for block in blocks:
            if block[0] == "paragraph":
                text, style_name = block[1], block[2]
                if not text:
                    continue
                is_heading = "heading" in style_name.lower() or "title" in style_name.lower()
                class_name = "preview-heading" if is_heading else "preview-paragraph"
                body_parts.append(f'<div class="{class_name}">{escape(text)}</div>')
                rendered_on_page += 1
            elif block[0] == "table":
                rows = block[1]
                body_parts.append('<table class="preview-table">')
                for row in rows:
                    body_parts.append('<tr>')
                    for cell in row:
                        body_parts.append(f'<td>{escape(_normalise_cell(cell))}</td>')
                    body_parts.append('</tr>')
                body_parts.append('</table>')
                rendered_on_page += max(1, len(rows))

            if rendered_on_page >= 28:
                body_parts.append(f'<div class="preview-page-number">Page {page_number}</div></div>')
                page_number += 1
                rendered_on_page = 0
                body_parts.append('<div class="preview-page">')

        body_parts.append(f'<div class="preview-page-number">Page {page_number}</div>')
    else:
        body_parts.append('<div class="preview-muted">No paragraph or table text found.</div>')

    body_parts.append('</div></div>')
    meta = f"{paragraph_count} paragraph(s) - {table_count} table(s) - no artificial preview limit"
    return _html_page("Microsoft Word", path.name, "".join(body_parts), meta)


def _preview_pptx_html(path: Path) -> str:
    """Render every PowerPoint slide as a slide-like card."""
    try:
        from pptx import Presentation
    except Exception:
        return _preview_pptx_html_from_zip(path)

    presentation = Presentation(str(path))
    body_parts = ['<div class="preview-document-pages">']
    for slide_number, slide in enumerate(presentation.slides, start=1):
        text_blocks = []
        for shape in sorted(slide.shapes, key=lambda shape: (getattr(shape, "top", 0), getattr(shape, "left", 0))):
            text = getattr(shape, "text", "")
            text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
            if text:
                text_blocks.append(text)

        body_parts.append('<div class="preview-slide">')
        body_parts.append(f'<div class="preview-slide-title">Slide {slide_number}</div>')
        if text_blocks:
            body_parts.append('<ul class="preview-list">')
            for text in text_blocks:
                body_parts.append(f'<li>{escape(text)}</li>')
            body_parts.append('</ul>')
        else:
            body_parts.append('<div class="preview-muted">No text found on this slide.</div>')
        body_parts.append('</div>')
    body_parts.append('</div>')
    return _html_page("Microsoft PowerPoint", path.name, "".join(body_parts), f"{len(presentation.slides)} slide(s) - no artificial preview limit")


def _preview_pptx_html_from_zip(path: Path) -> str:
    import re
    import zipfile
    import xml.etree.ElementTree as ET

    def slide_sort_key(name: str) -> int:
        match = re.search(r"slide(\d+)\.xml$", name)
        return int(match.group(1)) if match else 0

    with zipfile.ZipFile(path) as archive:
        slide_names = sorted(
            [name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")],
            key=slide_sort_key,
        )
        body_parts = ['<div class="preview-document-pages">']
        for slide_number, slide_name in enumerate(slide_names, start=1):
            root = ET.fromstring(archive.read(slide_name))
            text_nodes = []
            for node in root.iter():
                if node.tag.endswith("}t") and node.text:
                    value = node.text.strip()
                    if value:
                        text_nodes.append(value)

            body_parts.append('<div class="preview-slide">')
            body_parts.append(f'<div class="preview-slide-title">Slide {slide_number}</div>')
            if text_nodes:
                body_parts.append('<ul class="preview-list">')
                for value in text_nodes:
                    body_parts.append(f'<li>{escape(value)}</li>')
                body_parts.append('</ul>')
            else:
                body_parts.append('<div class="preview-muted">No text found on this slide.</div>')
            body_parts.append('</div>')
        body_parts.append('</div>')

    return _html_page("Microsoft PowerPoint", path.name, "".join(body_parts), f"{len(slide_names)} slide(s) - no artificial preview limit")


def _preview_xlsx_html(path: Path) -> str:
    """Render every workbook sheet and populated row without artificial limits."""
    try:
        from openpyxl import load_workbook
    except Exception as error:
        return _html_page("Microsoft Excel", path.name, f'<div class="preview-muted">Excel preview requires openpyxl. {escape(str(error))}</div>')

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    try:
        body_parts = ['<div class="preview-document-pages">']
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            body_parts.append('<div class="preview-sheet">')
            body_parts.append(f'<div class="preview-sheet-title">{escape(sheet.title)}</div>')
            rendered_any = False
            body_parts.append('<table class="preview-table">')
            for row in sheet.iter_rows(values_only=True):
                values = [_format_cell(value) for value in row]
                while values and values[-1] == "":
                    values.pop()
                if not values:
                    continue
                rendered_any = True
                body_parts.append('<tr>')
                for value in values:
                    body_parts.append(f'<td>{escape(value)}</td>')
                body_parts.append('</tr>')
            body_parts.append('</table>')
            if not rendered_any:
                body_parts.append('<div class="preview-muted">No cell values found.</div>')
            body_parts.append('</div>')
        body_parts.append('</div>')
        return _html_page("Microsoft Excel", path.name, "".join(body_parts), f"{len(workbook.sheetnames)} sheet(s) - no artificial preview limit")
    finally:
        workbook.close()


def _preview_docx_text(path: Path) -> str:
    try:
        from docx import Document
    except Exception as error:
        return f"DOCX preview requires python-docx.\n{error}"

    document = Document(str(path))
    lines = ["Microsoft Word document preview", ""]

    paragraphs = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    if paragraphs:
        lines.append("Paragraphs")
        lines.append("----------")
        lines.extend(paragraphs[:80])
        if len(paragraphs) > 80:
            lines.append(f"\n... showing first 80 paragraphs out of {len(paragraphs)}")
    else:
        lines.append("No paragraph text found.")

    if document.tables:
        lines.extend(["", "Tables", "------"])
        for table_index, table in enumerate(document.tables[:5], start=1):
            lines.append(f"Table {table_index}")
            rows = []
            for row in table.rows[:12]:
                cells = [_normalise_cell(cell.text) for cell in row.cells]
                rows.append(" | ".join(cells))
            lines.extend(rows or ["(empty table)"])
            if len(table.rows) > 12:
                lines.append(f"... showing first 12 rows out of {len(table.rows)}")
            lines.append("")
        if len(document.tables) > 5:
            lines.append(f"... showing first 5 tables out of {len(document.tables)}")

    return "\n".join(lines).strip()


def _preview_pptx_text(path: Path) -> str:
    import re
    import zipfile
    import xml.etree.ElementTree as ET

    def slide_sort_key(name: str) -> int:
        match = re.search(r"slide(\d+)\.xml$", name)
        return int(match.group(1)) if match else 0

    with zipfile.ZipFile(path) as archive:
        slide_names = sorted(
            [name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")],
            key=slide_sort_key,
        )
        lines = ["Microsoft PowerPoint preview", "", f"Slides: {len(slide_names)}"]

        for slide_number, slide_name in enumerate(slide_names[:20], start=1):
            root = ET.fromstring(archive.read(slide_name))
            text_nodes = []
            for node in root.iter():
                if node.tag.endswith("}t") and node.text:
                    value = node.text.strip()
                    if value:
                        text_nodes.append(value)

            lines.extend(["", f"Slide {slide_number}", "-------"])
            lines.extend(text_nodes or ["No text found on this slide."])

        if len(slide_names) > 20:
            lines.append(f"\n... showing first 20 slides out of {len(slide_names)}")

    return "\n".join(lines).strip()


def _preview_xlsx_text(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except Exception as error:
        return f"Excel preview requires openpyxl.\n{error}"

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    try:
        lines = ["Microsoft Excel workbook preview", "", f"Sheets: {len(workbook.sheetnames)}"]

        for sheet_name in workbook.sheetnames[:5]:
            sheet = workbook[sheet_name]
            lines.extend(["", f"Sheet: {sheet.title}", "-" * (7 + len(sheet.title))])
            rows = list(sheet.iter_rows(min_row=1, max_row=min(sheet.max_row or 1, 20), values_only=True))
            if not rows:
                lines.append("No cell values found.")
                continue

            rendered_rows = [_render_row(row) for row in rows if any(cell is not None for cell in row)]
            lines.extend(rendered_rows or ["No cell values found."])
            if sheet.max_row and sheet.max_row > 20:
                lines.append(f"... showing first 20 rows out of {sheet.max_row}")

        if len(workbook.sheetnames) > 5:
            lines.append(f"\n... showing first 5 sheets out of {len(workbook.sheetnames)}")

        return "\n".join(lines).strip()
    finally:
        workbook.close()


def _preview_archive_html(path: Path) -> str:
    title = path.name
    try:
        entries = _archive_entries(path)
    except Exception as error:
        body = f'<div class="preview-section"><div class="preview-muted">Could not inspect archive: {escape(str(error))}</div></div>'
        return _html_page("Archive", title, body)

    if not entries:
        body = '<div class="preview-section"><div class="preview-muted">No entries found in this archive.</div></div>'
        return _html_page("Archive", title, body, "0 item(s)")

    body_parts = ['<div class="preview-section"><div class="preview-section-title">Archive contents</div>']
    body_parts.append('<table class="preview-table"><tr><th>Name</th><th>Size</th></tr>')
    for name, size in entries[:80]:
        body_parts.append(f'<tr><td>{escape(name)}</td><td>{escape(_format_bytes(size))}</td></tr>')
    body_parts.append('</table>')
    if len(entries) > 80:
        body_parts.append(f'<div class="preview-muted">Showing first 80 entries out of {len(entries)}.</div>')
    body_parts.append('</div>')
    return _html_page("Archive", title, "".join(body_parts), f"{len(entries)} item(s)")


def _preview_archive_text(path: Path) -> str:
    try:
        entries = _archive_entries(path)
    except Exception as error:
        return f"Archive preview unavailable\n\nCould not inspect archive: {error}"

    lines = ["Archive preview", "", f"Items: {len(entries)}", ""]
    for name, size in entries[:80]:
        lines.append(f"{name}    {_format_bytes(size)}")
    if len(entries) > 80:
        lines.append(f"\n... showing first 80 entries out of {len(entries)}")
    return "\n".join(lines).strip()


def _archive_entries(path: Path):
    path = Path(path)
    suffix = _preview_suffix(path)

    if suffix == ".zip":
        with zipfile.ZipFile(path) as archive:
            return [(info.filename, info.file_size) for info in archive.infolist()]

    if suffix in {".tar", ".tgz", ".tar.gz"}:
        with tarfile.open(path) as archive:
            return [(member.name, member.size) for member in archive.getmembers()]

    if suffix == ".gz":
        with gzip.open(path, "rb") as archive:
            sample = archive.read(1024)
        return [(path.with_suffix("").name or path.name, len(sample))]

    raise ValueError("This archive type cannot be inspected without an external parser.")


def _iter_docx_blocks(document):
    """Yield paragraphs and tables in document order for a less scrambled DOCX preview."""
    try:
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except Exception:
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                yield ("paragraph", text, getattr(paragraph.style, "name", ""))
        for table in document.tables:
            yield ("table", [[cell.text for cell in row.cells] for row in table.rows])
        return

    body = document.element.body
    for child in body.iterchildren():
        if isinstance(child, CT_P):
            paragraph = Paragraph(child, document)
            text = paragraph.text.strip()
            if text:
                yield ("paragraph", text, getattr(paragraph.style, "name", ""))
        elif isinstance(child, CT_Tbl):
            table = Table(child, document)
            yield ("table", [[cell.text for cell in row.cells] for row in table.rows])


def _render_row(row: Iterable[object]) -> str:
    values = [_format_cell(value) for value in row]
    while values and values[-1] == "":
        values.pop()
    return " | ".join(values)


def _format_cell(value: object) -> str:
    if value is None:
        return ""
    return str(value).replace("\n", " ").strip()


def _format_bytes(value) -> str:
    try:
        size = int(value or 0)
    except Exception:
        return "Unknown"
    units = ["B", "KB", "MB", "GB", "TB"]
    amount = float(size)
    for unit in units:
        if amount < 1024 or unit == units[-1]:
            return f"{amount:.1f} {unit}" if unit != "B" else f"{int(amount)} B"
        amount /= 1024


def _normalise_cell(text: str) -> str:
    return " ".join((text or "").split())


def _truncate(text: str, max_chars: int | None) -> str:
    text = text or ""
    if max_chars is None:
        return text
    if len(text) > max_chars:
        return text[:max_chars] + "\n\n... preview truncated"
    return text


def _truncate_html(text: str, max_chars: int | None) -> str:
    text = text or ""
    if max_chars is None:
        return text
    if len(text) > max_chars:
        return text[:max_chars] + '<div class="preview-muted">Preview truncated.</div></div>'
    return text

from __future__ import annotations

from pathlib import Path
from typing import Literal

from core.helpers import unique_path

MicrosoftFileKind = Literal["document", "powerpoint", "excel"]


class MicrosoftSuiteError(RuntimeError):
    """Raised when a Microsoft file cannot be created."""


MICROSOFT_EXTENSIONS = {
    "document": ".docx",
    "powerpoint": ".pptx",
    "excel": ".xlsx",
}

MICROSOFT_TITLES = {
    "document": "Microsoft Word Document",
    "powerpoint": "Microsoft PowerPoint Presentation",
    "excel": "Microsoft Excel Workbook",
}


def normalise_office_name(name: str, kind: MicrosoftFileKind) -> str:
    base = (name or "").strip() or MICROSOFT_TITLES[kind]
    wanted_suffix = MICROSOFT_EXTENSIONS[kind]
    if not Path(base).suffix:
        base = f"{base}{wanted_suffix}"
    return base


def create_microsoft_file(destination_dir: str | Path, name: str, kind: MicrosoftFileKind) -> Path:
    """Create a real local Microsoft Office file.

    Heavy lifting is delegated to the standard Python Office libraries when
    available. Imports are intentionally local so the app can still launch even
    before optional dependencies are installed.
    """
    if kind not in MICROSOFT_EXTENSIONS:
        raise MicrosoftSuiteError(f"Unsupported Microsoft file kind: {kind}")

    destination_dir = Path(destination_dir)
    destination_dir.mkdir(parents=True, exist_ok=True)
    filename = normalise_office_name(name, kind)
    target = unique_path(destination_dir, filename)

    try:
        if kind == "document":
            from docx import Document  # type: ignore

            doc = Document()
            doc.add_heading(Path(target).stem, level=1)
            doc.add_paragraph("Created from ZJX-LMS.")
            doc.save(str(target))
            return target

        if kind == "powerpoint":
            from pptx import Presentation  # type: ignore

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            slide.shapes.title.text = Path(target).stem
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "Created from ZJX-LMS."
            presentation.save(str(target))
            return target

        if kind == "excel":
            from openpyxl import Workbook  # type: ignore

            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "Sheet1"
            sheet["A1"] = "Created from ZJX-LMS"
            workbook.save(str(target))
            return target
    except ImportError as exc:
        raise MicrosoftSuiteError(
            "Missing Microsoft creation dependency. Install python-docx, python-pptx, and openpyxl."
        ) from exc
    except Exception as exc:
        raise MicrosoftSuiteError(str(exc)) from exc

    raise MicrosoftSuiteError(f"Unsupported Microsoft file kind: {kind}")
